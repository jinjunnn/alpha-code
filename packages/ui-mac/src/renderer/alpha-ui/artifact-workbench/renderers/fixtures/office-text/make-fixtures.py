#!/usr/bin/env python3
"""Regenerate the office-text fixtures from real generators (REQ-123 / #1175).

Requires: macOS `textutil`, python-docx, python-pptx. Writes the extracted part
files into this directory, mirroring the real in-package part names. See README.md
for why assertions must be content-string based, never structure based.
"""

import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path

HERE = Path(__file__).resolve().parent

COCOA_RTF = rb"""{\rtf1\ansi\deff0
{\fonttbl{\f0 Helvetica;}}
\f0\fs48\b Cocoa Heading Line\b0\fs24\par
Cocoa intro paragraph text with plain body words.\par
\bullet\tab cocoa bullet one\par
\bullet\tab cocoa bullet two\par
\trowd\cellx3000\cellx6000 CocoaCell Widget\intbl\cell CocoaCell Qty\intbl\cell\row
\trowd\cellx3000\cellx6000 CocoaCell Anvil\intbl\cell CocoaCell 42\intbl\cell\row
}
"""


def extract(container: Path, parts: dict[str, Path]) -> None:
    with zipfile.ZipFile(container) as zf:
        for part_name, dest in parts.items():
            dest.parent.mkdir(parents=True, exist_ok=True)
            dest.write_bytes(zf.read(part_name))


def make_cocoa_docx(tmp: Path) -> None:
    rtf = tmp / "cocoa-src.rtf"
    rtf.write_bytes(COCOA_RTF)
    out = tmp / "cocoa.docx"
    subprocess.run(
        ["textutil", "-convert", "docx", str(rtf), "-output", str(out)], check=True
    )
    extract(out, {"word/document.xml": HERE / "cocoa-docx/word/document.xml"})


def make_py_docx(tmp: Path) -> None:
    from docx import Document

    d = Document()
    d.add_heading("Quarterly Report Heading", level=1)
    p = d.add_paragraph("Intro paragraph with ")
    p.add_run("bold emphasis").bold = True
    d.add_paragraph("first bullet item", style="List Bullet")
    d.add_paragraph("second bullet item", style="List Bullet")
    tbl = d.add_table(rows=2, cols=2)
    tbl.cell(0, 0).text = "Widget"
    tbl.cell(0, 1).text = "Qty"
    tbl.cell(1, 0).text = "Anvil"
    tbl.cell(1, 1).text = "42"
    out = tmp / "pydocx.docx"
    d.save(out)
    extract(out, {"word/document.xml": HERE / "py-docx/word/document.xml"})


def make_py_pptx(tmp: Path) -> None:
    from pptx import Presentation

    prs = Presentation()

    def add(title: str, body: str, notes: str | None = None) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
        if notes is not None:
            slide.notes_slide.notes_text_frame.text = notes

    add("Alpha Slide One", "bullet one alpha", "note for alpha")
    add("Bravo Slide Two", "bullet two bravo")  # deliberately no notes part
    add("Charlie Slide Three", "bullet three charlie", "note for charlie")

    # Real user operation: move the third slide to the front, so the authoritative
    # sldIdLst order differs from the filename order (baseline ②-4).
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001 — supported escape hatch
    ids = list(sld_id_lst)
    sld_id_lst.remove(ids[2])
    sld_id_lst.insert(0, ids[2])
    out = tmp / "py.pptx"
    prs.save(out)

    extract(
        out,
        {
            "ppt/presentation.xml": HERE / "py-pptx/ppt/presentation.xml",
            "ppt/_rels/presentation.xml.rels": HERE / "py-pptx/ppt/_rels/presentation.xml.rels",
            "ppt/slides/slide1.xml": HERE / "py-pptx/ppt/slides/slide1.xml",
            "ppt/slides/slide2.xml": HERE / "py-pptx/ppt/slides/slide2.xml",
            "ppt/slides/slide3.xml": HERE / "py-pptx/ppt/slides/slide3.xml",
            "ppt/slides/_rels/slide1.xml.rels": HERE / "py-pptx/ppt/slides/_rels/slide1.xml.rels",
            "ppt/slides/_rels/slide2.xml.rels": HERE / "py-pptx/ppt/slides/_rels/slide2.xml.rels",
            "ppt/slides/_rels/slide3.xml.rels": HERE / "py-pptx/ppt/slides/_rels/slide3.xml.rels",
            "ppt/notesSlides/notesSlide1.xml": HERE / "py-pptx/ppt/notesSlides/notesSlide1.xml",
            "ppt/notesSlides/notesSlide2.xml": HERE / "py-pptx/ppt/notesSlides/notesSlide2.xml",
        },
    )


def main() -> None:
    tmpdir = Path(tempfile.mkdtemp(prefix="office-text-fixtures-"))
    try:
        make_cocoa_docx(tmpdir)
        make_py_docx(tmpdir)
        make_py_pptx(tmpdir)
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)
    print("fixtures regenerated under", HERE)


if __name__ == "__main__":
    main()
