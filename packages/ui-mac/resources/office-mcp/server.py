#!/usr/bin/env python3
"""Alpha first-party Office MCP server (stdio only).

One bundled server exposes a narrow tool pair for the selected file format. The
host supplies both the format and a canonical workspace directory on argv; tool
callers can only pass absolute paths below that directory.
"""

from __future__ import annotations

import json
import os
from pathlib import Path
import sys
import tempfile
from typing import Any, Callable


PROTOCOL_VERSION = "2025-06-18"
SERVER_VERSION = "1.0.0"
FORMATS = {"word", "excel", "powerpoint", "pdf"}
EXTENSIONS = {"word": ".docx", "excel": ".xlsx", "powerpoint": ".pptx", "pdf": ".pdf"}


def main(argv: list[str]) -> int:
    if len(argv) != 3 or argv[1] not in FORMATS:
        print("usage: server.py <word|excel|powerpoint|pdf> <absolute-workspace>; stdio is the only transport", file=sys.stderr)
        return 2

    workspace_arg = argv[2]
    if not is_absolute_path(workspace_arg) or has_traversal(workspace_arg):
        print("workspace must be an absolute path without '..' segments", file=sys.stderr)
        return 2

    try:
        workspace = Path(workspace_arg).resolve(strict=True)
    except OSError:
        print("workspace must be an existing directory", file=sys.stderr)
        return 2
    if not workspace.is_dir():
        print("workspace must be an existing directory", file=sys.stderr)
        return 2

    for line in sys.stdin:
        if not line.strip():
            continue
        message: Any = None
        try:
            message = json.loads(line)
            response = handle_message(argv[1], workspace, message)
        except Exception as error:
            request_id = message.get("id") if isinstance(message, dict) else None
            response = error_response(request_id, -32603, str(error))
        if response is not None:
            print(json.dumps(response, ensure_ascii=False, default=str), flush=True)
    return 0


def handle_message(format_name: str, workspace: Path, message: Any) -> dict[str, Any] | None:
    if not isinstance(message, dict) or message.get("jsonrpc") != "2.0":
        return error_response(message.get("id") if isinstance(message, dict) else None, -32600, "invalid JSON-RPC request")

    method = message.get("method")
    request_id = message.get("id")
    if method == "notifications/initialized":
        return None
    if method == "initialize":
        return success_response(
            request_id,
            {
                "protocolVersion": PROTOCOL_VERSION,
                "capabilities": {"tools": {"listChanged": False}},
                "serverInfo": {"name": f"alpha-office-{format_name}", "version": SERVER_VERSION},
            },
        )
    if method == "ping":
        return success_response(request_id, {})
    if method == "tools/list":
        return success_response(request_id, {"tools": tools_for(format_name)})
    if method == "tools/call":
        params = message.get("params")
        if not isinstance(params, dict) or not isinstance(params.get("name"), str):
            return error_response(request_id, -32602, "tools/call requires a tool name")
        arguments = params.get("arguments", {})
        if not isinstance(arguments, dict):
            return error_response(request_id, -32602, "tool arguments must be an object")
        try:
            result = call_tool(format_name, workspace, params["name"], arguments)
            return success_response(
                request_id,
                {"content": [{"type": "text", "text": json.dumps(result, ensure_ascii=False, default=str)}]},
            )
        except Exception as error:
            return success_response(
                request_id,
                {"content": [{"type": "text", "text": str(error)}], "isError": True},
            )
    if request_id is None:
        return None
    return error_response(request_id, -32601, f"method not found: {method}")


def tools_for(format_name: str) -> list[dict[str, Any]]:
    path_property = {"path": {"type": "string", "description": "Absolute path inside the granted workspace"}}
    if format_name == "word":
        return [
            tool("read_docx", "Read paragraphs and tables from a Word document", path_property, ["path"]),
            tool(
                "write_docx",
                "Create, replace, or append text to a Word document",
                {
                    **path_property,
                    "title": {"type": "string"},
                    "paragraphs": {"type": "array", "items": {"type": "string"}},
                    "append": {"type": "boolean", "default": False},
                },
                ["path", "paragraphs"],
            ),
        ]
    if format_name == "excel":
        return [
            tool("read_xlsx", "Read cell values from every worksheet in an Excel workbook", path_property, ["path"]),
            tool(
                "write_xlsx",
                "Create a workbook or update addressed cells in an existing workbook",
                {
                    **path_property,
                    "sheets": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "name": {"type": "string"},
                                "cells": {"type": "object", "additionalProperties": True},
                            },
                            "required": ["name", "cells"],
                            "additionalProperties": False,
                        },
                    },
                },
                ["path", "sheets"],
            ),
        ]
    if format_name == "powerpoint":
        return [
            tool("read_pptx", "Read text from every slide in a PowerPoint deck", path_property, ["path"]),
            tool(
                "write_pptx",
                "Create, replace, or append text slides in a PowerPoint deck",
                {
                    **path_property,
                    "slides": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string"},
                                "body": {"oneOf": [{"type": "string"}, {"type": "array", "items": {"type": "string"}}]},
                            },
                            "required": ["title", "body"],
                            "additionalProperties": False,
                        },
                    },
                    "append": {"type": "boolean", "default": False},
                },
                ["path", "slides"],
            ),
        ]
    return [
        tool("read_pdf", "Extract the text layer from every PDF page", path_property, ["path"]),
        tool(
            "write_pdf",
            "Generate a text PDF, replace its text document, or append text pages",
            {
                **path_property,
                "pages": {"type": "array", "items": {"type": "string"}},
                "mode": {"type": "string", "enum": ["replace", "append"], "default": "replace"},
            },
            ["path", "pages"],
        ),
    ]


def tool(name: str, description: str, properties: dict[str, Any], required: list[str]) -> dict[str, Any]:
    return {
        "name": name,
        "description": description,
        "inputSchema": {
            "type": "object",
            "properties": properties,
            "required": required,
            "additionalProperties": False,
        },
    }


def call_tool(format_name: str, workspace: Path, name: str, arguments: dict[str, Any]) -> dict[str, Any]:
    handlers: dict[str, dict[str, Callable[[Path, dict[str, Any]], dict[str, Any]]]] = {
        "word": {"read_docx": read_docx, "write_docx": write_docx},
        "excel": {"read_xlsx": read_xlsx, "write_xlsx": write_xlsx},
        "powerpoint": {"read_pptx": read_pptx, "write_pptx": write_pptx},
        "pdf": {"read_pdf": read_pdf, "write_pdf": write_pdf},
    }
    handler = handlers[format_name].get(name)
    if handler is None:
        raise ValueError(f"unknown {format_name} tool: {name}")
    path = require_workspace_path(workspace, arguments.get("path"), EXTENSIONS[format_name])
    return handler(path, arguments)


def require_workspace_path(workspace: Path, raw: Any, extension: str) -> Path:
    if not isinstance(raw, str) or not raw:
        raise ValueError("path must be a non-empty string")
    if not is_absolute_path(raw):
        raise ValueError("path must be absolute")
    if has_traversal(raw):
        raise ValueError("path traversal ('..') is not allowed")
    candidate = Path(raw)
    if candidate.suffix.lower() != extension:
        raise ValueError(f"path must end in {extension}")
    resolved = candidate.resolve(strict=False)
    if os.path.commonpath((str(workspace), str(resolved))) != str(workspace):
        raise ValueError("path is outside the granted workspace")
    return resolved


def is_absolute_path(value: str) -> bool:
    normalized = value.replace("\\", "/")
    return normalized.startswith("/") or (len(normalized) >= 3 and normalized[0].isalpha() and normalized[1:3] == ":/")


def has_traversal(value: str) -> bool:
    return ".." in value.replace("\\", "/").split("/")


def read_docx(path: Path, _: dict[str, Any]) -> dict[str, Any]:
    require_existing(path)
    from docx import Document

    document = Document(path)
    return {
        "path": str(path),
        "paragraphs": [paragraph.text for paragraph in document.paragraphs],
        "tables": [
            [[cell.text for cell in row.cells] for row in table.rows]
            for table in document.tables
        ],
    }


def write_docx(path: Path, arguments: dict[str, Any]) -> dict[str, Any]:
    from docx import Document

    paragraphs = require_string_list(arguments.get("paragraphs"), "paragraphs")
    append = arguments.get("append", False)
    if not isinstance(append, bool):
        raise ValueError("append must be a boolean")
    document = Document(path) if append and path.exists() else Document()
    title = arguments.get("title")
    if title is not None:
        if not isinstance(title, str):
            raise ValueError("title must be a string")
        document.add_heading(title, level=0)
    for paragraph in paragraphs:
        document.add_paragraph(paragraph)
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)
    return {"path": str(path), "paragraphsWritten": len(paragraphs), "appended": append}


def read_xlsx(path: Path, _: dict[str, Any]) -> dict[str, Any]:
    require_existing(path)
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=False)
    return {
        "path": str(path),
        "sheets": [
            {"name": sheet.title, "rows": [list(row) for row in sheet.iter_rows(values_only=True)]}
            for sheet in workbook.worksheets
        ],
    }


def write_xlsx(path: Path, arguments: dict[str, Any]) -> dict[str, Any]:
    from openpyxl import Workbook, load_workbook

    sheets = arguments.get("sheets")
    if not isinstance(sheets, list) or not sheets:
        raise ValueError("sheets must be a non-empty array")
    workbook = load_workbook(path) if path.exists() else Workbook()
    created = not path.exists()
    touched: list[str] = []
    for item in sheets:
        if not isinstance(item, dict) or not isinstance(item.get("name"), str) or not item["name"]:
            raise ValueError("each sheet requires a non-empty name")
        cells = item.get("cells")
        if not isinstance(cells, dict):
            raise ValueError("each sheet requires a cells object")
        sheet = workbook[item["name"]] if item["name"] in workbook.sheetnames else workbook.create_sheet(item["name"])
        for coordinate, value in cells.items():
            if not isinstance(coordinate, str):
                raise ValueError("cell addresses must be strings")
            sheet[coordinate] = value
        touched.append(item["name"])
    if created and "Sheet" in workbook.sheetnames and "Sheet" not in touched and len(workbook.sheetnames) > 1:
        workbook.remove(workbook["Sheet"])
    path.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(path)
    return {"path": str(path), "sheetsUpdated": touched}


def read_pptx(path: Path, _: dict[str, Any]) -> dict[str, Any]:
    require_existing(path)
    from pptx import Presentation

    presentation = Presentation(path)
    return {
        "path": str(path),
        "slides": [
            {
                "number": index,
                "text": [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text],
            }
            for index, slide in enumerate(presentation.slides, start=1)
        ],
    }


def write_pptx(path: Path, arguments: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation

    slides = arguments.get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("slides must be a non-empty array")
    append = arguments.get("append", False)
    if not isinstance(append, bool):
        raise ValueError("append must be a boolean")
    presentation = Presentation(path) if append and path.exists() else Presentation()
    for item in slides:
        if not isinstance(item, dict) or not isinstance(item.get("title"), str):
            raise ValueError("each slide requires a string title")
        body = item.get("body")
        if isinstance(body, list):
            body = "\n".join(require_string_list(body, "slide body"))
        if not isinstance(body, str):
            raise ValueError("each slide body must be a string or string array")
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = item["title"]
        slide.placeholders[1].text = body
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)
    return {"path": str(path), "slidesWritten": len(slides), "appended": append}


def read_pdf(path: Path, _: dict[str, Any]) -> dict[str, Any]:
    require_existing(path)
    from pypdf import PdfReader

    reader = PdfReader(path)
    return {"path": str(path), "pages": [page.extract_text() or "" for page in reader.pages]}


def write_pdf(path: Path, arguments: dict[str, Any]) -> dict[str, Any]:
    from pypdf import PdfWriter

    pages = require_string_list(arguments.get("pages"), "pages")
    if not pages:
        raise ValueError("pages must be non-empty")
    mode = arguments.get("mode", "replace")
    if mode not in {"replace", "append"}:
        raise ValueError("mode must be replace or append")
    if mode == "append":
        require_existing(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(suffix=".pdf", dir=path.parent, delete=False) as temporary:
        generated = Path(temporary.name)
    merged: Path | None = None
    try:
        render_text_pdf(generated, pages)
        if mode == "replace":
            os.replace(generated, path)
        else:
            writer = PdfWriter()
            writer.append(path)
            writer.append(generated)
            with tempfile.NamedTemporaryFile(suffix=".pdf", dir=path.parent, delete=False) as output:
                merged = Path(output.name)
                writer.write(output)
            os.replace(merged, path)
    finally:
        generated.unlink(missing_ok=True)
        if merged is not None:
            merged.unlink(missing_ok=True)
    return {"path": str(path), "pagesWritten": len(pages), "mode": mode}


def render_text_pdf(path: Path, pages: list[str]) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    pdf = canvas.Canvas(str(path), pagesize=A4)
    _, height = A4
    for page in pages:
        text = pdf.beginText(54, height - 54)
        text.setFont("Helvetica", 11)
        for source_line in page.splitlines() or [""]:
            chunks = [source_line[index : index + 95] for index in range(0, len(source_line), 95)] or [""]
            for chunk in chunks:
                if text.getY() < 54:
                    pdf.drawText(text)
                    pdf.showPage()
                    text = pdf.beginText(54, height - 54)
                    text.setFont("Helvetica", 11)
                text.textLine(chunk)
        pdf.drawText(text)
        pdf.showPage()
    pdf.save()


def require_string_list(value: Any, name: str) -> list[str]:
    if not isinstance(value, list) or any(not isinstance(item, str) for item in value):
        raise ValueError(f"{name} must be an array of strings")
    return value


def require_existing(path: Path) -> None:
    if not path.is_file():
        raise ValueError(f"file does not exist: {path}")


def success_response(request_id: Any, result: Any) -> dict[str, Any]:
    return {"jsonrpc": "2.0", "id": request_id, "result": result}


def error_response(request_id: Any, code: int, message: str) -> dict[str, Any]:
    return {"jsonrpc": "2.0", "id": request_id, "error": {"code": code, "message": message}}


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
